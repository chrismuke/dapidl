"""Assemble the A-pair + annotation deck (.pptx) for the lab meeting.

Order:
  1. Title
  2. Pseudo-GT sources (fig_a0)
  3. Pipeline overview (fig_a6)
  4. Headline: granularity × cross-source (fig_a1)
  5. Per-class breakdown heatmap (fig_a2)
  6. Annotation method baseline (fig_a3)
  7. DAPI vs annotation per-class (fig_a4)
  8. Training dynamics (fig_a5)
  9. Take-aways
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FIG_DIR = Path("/mnt/work/git/dapidl/pipeline_output/figures_v2")
OUT = Path("/mnt/work/git/dapidl/presentation/DAPIDL_a_pair_2026_05_04.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x1A, 0x4B, 0x7A)
BODY = RGBColor(0x2B, 0x2B, 0x2B)
MUTED = RGBColor(0x6E, 0x6E, 0x6E)
ACCENT = RGBColor(0xD4, 0xA0, 0x17)


def add_text(slide, text: str, left: float, top: float, width: float, height: float,
             size: int = 16, color: RGBColor = BODY, bold: bool = False,
             align=PP_ALIGN.LEFT, italic: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_figure_slide(prs, title: str, fig_path: Path, takeaway: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_text(slide, title, 0.6, 0.35, 12.1, 0.7, size=26, color=NAVY, bold=True)
    # Figure (full-width, leaves room for title + takeaway)
    if fig_path.exists():
        slide.shapes.add_picture(
            str(fig_path), Inches(0.6), Inches(1.15),
            width=Inches(12.1),
        )
    # Takeaway at bottom
    add_text(slide, "▶ " + takeaway, 0.6, 6.95, 12.1, 0.4,
             size=12, color=BODY, italic=True)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide,
             "DAPIDL — Cell-type prediction from DAPI nuclear staining",
             0.6, 2.5, 12.1, 1.2, size=36, color=NAVY, bold=True)
    add_text(slide,
             "First clean cross-source numbers + annotation baseline",
             0.6, 3.6, 12.1, 0.6, size=20, color=BODY, italic=True)
    add_text(slide,
             "Lab meeting · 2026-05-04 · A pair (Janesick rep1+rep2 → STHELAR breast s0/s1/s3/s6)",
             0.6, 6.5, 12.1, 0.5, size=11, color=MUTED, italic=True)


def add_takeaways_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, "Take-aways", 0.6, 0.4, 12.1, 0.7,
             size=30, color=NAVY, bold=True)

    bullets = [
        ("Granularity matters",
         "COARSE 4-class works (val F1 0.753); MEDIUM 12-class drops sharply (0.497)."),
        ("Cross-source gap is real",
         "Janesick → STHELAR drops F1 by ~37 pp at COARSE, ~34 pp at MEDIUM."),
        ("Epithelial transfers; rare classes don't",
         "Mammary luminal F1 0.65-0.70 across all slides. Adipocyte/Pericyte → 0."),
        ("DAPI is competitive on common compartments",
         "Per-class DAPI ≈ scType on Epi/Imm/Stromal. Endothelial hard for everyone."),
        ("BANKSY+scType is the best annotation baseline",
         "Mean F1 0.568 on COARSE-4 (beats scType 0.558). Spatial-aware clustering helps; SingleR still drags ensembles."),
        ("Multi-source pooling beats single-source training",
         "STHELAR-Prime alone transfers 0.23 to Janesick rep1. Pooling all STHELAR sources lifts this to 0.62 (~3×)."),
        ("Confidence: STHELAR cells_label2 is solid",
         "Median per-cell confidence = 1.000 (q05 0.764). Equivalent to cells_final_label at COARSE."),
        ("What's next",
         "Skin tissue (Coarse F1 0.765, Medium 0.582) ready to fold in. CelloType instance-seg + breast_s1 BANKSY workaround deferred."),
    ]
    y = 1.3
    for title, sub in bullets:
        add_text(slide, "● " + title, 0.7, y, 12.0, 0.4,
                 size=16, color=NAVY, bold=True)
        add_text(slide, "    " + sub, 0.7, y + 0.4, 12.0, 0.45,
                 size=13, color=BODY)
        y += 0.85


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_figure_slide(prs, "Pseudo-ground-truth sources",
                     FIG_DIR / "fig_a0_data_pseudo_gt.png",
                     "Janesick provides expert 17-class GT (gold). STHELAR provides cells_label2 — a multi-method consensus with median per-cell confidence 1.0.")

    add_figure_slide(prs, "Pipeline overview",
                     FIG_DIR / "fig_a6_pipeline.png",
                     "Both label sources feed 128px DAPI patches into EfficientNetV2-S → COARSE-4 / MEDIUM-12 prediction.")

    add_figure_slide(prs, "Granularity × cross-source generalization",
                     FIG_DIR / "fig_a1_headline.png",
                     "DAPI predicts COARSE-4 well in-domain (0.753); cross-source domain shift costs ~37 pp. MEDIUM-12 collapses to 0.13-0.19 on STHELAR.")

    add_figure_slide(prs, "Cross-source matrix — A/B/C/D × per-slide test",
                     FIG_DIR / "fig_a7_cross_source_matrix.png",
                     "Pooling heterogeneous training data (D) turns Prime-only's 0.23 → 0.62 transfer on Janesick rep1 — a ~3× improvement at COARSE-4.")

    add_figure_slide(prs, "Per-class breakdown — what transfers and what doesn't",
                     FIG_DIR / "fig_a2_per_class_heatmap.png",
                     "Mammary Luminal transfers cleanly (F1 0.65-0.70). Rare types (Adipocyte, Pericyte, Dendritic_Cell) collapse to 0 — class imbalance + inherent difficulty.")

    add_figure_slide(prs, "Annotation method baseline (validates our labels)",
                     FIG_DIR / "fig_a3_annotation.png",
                     "scType custom_default beats CellTypist + SingleR on every slide. CT+scType consensus is best 2-method (0.538); adding SingleR HURTS the ensemble.")

    add_figure_slide(prs, "Per-class DAPI vs gene-expression baselines",
                     FIG_DIR / "fig_a4_class_profile.png",
                     "DAPI is competitive on the dominant classes (Epithelial, Immune, Stromal). Endothelial is the hardest class for every method — biological signal is too sparse.")

    add_figure_slide(prs, "Training dynamics — how the model converged",
                     FIG_DIR / "fig_a5_training_curves.png",
                     "Coarse: best ep 12 (val F1 0.753), early-stopped ep 20. Medium: best ep 14 (val F1 0.497), early-stopped ep 22. Both converge cleanly.")

    add_takeaways_slide(prs)

    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()

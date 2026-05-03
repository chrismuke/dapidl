"""Build the DAPIDL v2 deck (.pptx) implementing SLIDE_REWRITE_V2.md.

Design language:
- 16:9 widescreen (13.33" × 7.5")
- White backgrounds, navy titles, generous white space
- One idea per slide, consistent typography, captions in 11pt grey
- Figure slides: full-width figure + 1-line takeaway

Color contract matches scripts/figures_v2/_style.py.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


FIG_DIR = Path("/mnt/work/git/dapidl/pipeline_output/figures_v2")
EXIST_DIR = Path("/mnt/work/git/dapidl/pipeline_output")
OUT = Path("/mnt/work/git/dapidl/presentation/DAPIDL_v2.pptx")


# ─── Design tokens ─────────────────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x4B, 0x7A)       # primary title / accent
NAVY_DARK = RGBColor(0x10, 0x35, 0x5A)  # hover / strong
ACCENT = RGBColor(0xB8, 0xDA, 0xEC)     # light blue chips
BODY = RGBColor(0x2B, 0x2B, 0x2B)       # body text
MUTED = RGBColor(0x6E, 0x6E, 0x6E)      # captions
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EPI_RED = RGBColor(0xED, 0x55, 0x3B)
IMM_NAVY = RGBColor(0x3D, 0x5A, 0x80)
STR_BROWN = RGBColor(0xBC, 0x8D, 0x5A)
END_TEAL = RGBColor(0x4A, 0x9D, 0x9C)


# ─── Helpers ───────────────────────────────────────────────────────────────

def add_title(slide, text: str, top_in: float = 0.4, height_in: float = 0.9,
              size_pt: int = 32, color: RGBColor = NAVY,
              left_in: float = 0.6, width_in: float = 12.1) -> None:
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.bold = True
    run.font.color.rgb = color


def add_subtitle(slide, text: str, top_in: float = 1.25,
                 size_pt: int = 16, color: RGBColor = MUTED,
                 left_in: float = 0.6, width_in: float = 12.1,
                 italic: bool = True) -> None:
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.italic = italic
    run.font.color.rgb = color


def add_paragraph_block(slide, lines: list[str], left_in: float, top_in: float,
                        width_in: float, height_in: float,
                        size_pt: int = 16, color: RGBColor = BODY,
                        bullet: bool = True, line_spacing: float = 1.2) -> None:
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ("•  " + line) if bullet else line
        run.font.name = "Calibri"
        run.font.size = Pt(size_pt)
        run.font.color.rgb = color


def add_caption(slide, text: str, top_in: float = 7.05,
                left_in: float = 0.6, width_in: float = 12.1,
                size_pt: int = 10) -> None:
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.italic = True
    run.font.color.rgb = MUTED


def add_image_centered(slide, path: Path,
                       top_in: float = 1.6, max_h_in: float = 5.2,
                       max_w_in: float = 12.0) -> None:
    """Place an image centered on the slide, scaled to fit the box."""
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size
    aspect = iw / ih
    # Try to fit by width first; fall back to height
    w_in = max_w_in
    h_in = w_in / aspect
    if h_in > max_h_in:
        h_in = max_h_in
        w_in = h_in * aspect
    left_in = (13.333 - w_in) / 2
    slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in),
                             width=Inches(w_in), height=Inches(h_in))


def add_left_accent_bar(slide, top_in: float = 0.45, height_in: float = 0.85,
                        width_in: float = 0.10) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.45), Inches(top_in),
                                 Inches(width_in), Inches(height_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()


def add_box(slide, text_lines: list[tuple[str, int, RGBColor, bool]],
            left_in: float, top_in: float, width_in: float, height_in: float,
            fill: RGBColor = ACCENT, border: RGBColor | None = None) -> None:
    """A filled box with stacked text. Each line: (text, size_pt, color, bold)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, size, col, bold) in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col


# ─── Slide builders ────────────────────────────────────────────────────────

def slide_section_header(prs: Presentation, kicker: str, title: str,
                         takeaway: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Kicker (small caps colored text)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = kicker.upper()
    r.font.name = "Calibri"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = NAVY
    # Title
    add_title(slide, title, top_in=0.85, height_in=1.1, size_pt=42)
    # Takeaway
    add_subtitle(slide, takeaway, top_in=2.1, size_pt=20, italic=True)
    return slide


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Date in corner
    box = slide.shapes.add_textbox(Inches(11.2), Inches(0.25),
                                   Inches(2.0), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "v2 · 2026-05-02"
    r.font.name = "Calibri"
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED

    # Massive title
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.3),
                                   Inches(12.1), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "DAPIDL"
    r.font.name = "Calibri"
    r.font.size = Pt(110)
    r.font.bold = True
    r.font.color.rgb = NAVY
    # Subtitle
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = "Predicting cell types from DAPI nuclei"
    r.font.name = "Calibri"
    r.font.size = Pt(28)
    r.font.color.rgb = BODY
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run()
    r.text = "Deep learning + spatial transcriptomics"
    r.font.name = "Calibri"
    r.font.size = Pt(20)
    r.font.italic = True
    r.font.color.rgb = MUTED

    # Footer line
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.7),
                                   Inches(12.1), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Christian Muke   ·   v2 deck rebuilt around measurements"
    r.font.name = "Calibri"
    r.font.size = Pt(13)
    r.font.color.rgb = MUTED


def slide_goal(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_accent_bar(slide)
    add_title(slide, "Goal", left_in=0.75)
    # Light blue box with the one-line goal
    add_box(slide, [
        ("Automatic cell-type annotation from DAPI-stained nuclei alone — "
         "trained on labels from spatial transcriptomics.", 18, NAVY_DARK, True),
    ], left_in=0.6, top_in=1.55, width_in=12.1, height_in=0.95)

    add_paragraph_block(slide, [
        "Train a neural network to predict cell types from nuclear morphology",
        "Use spatial transcriptomics (Xenium, MERSCOPE, STHELAR) to generate labels at scale",
        "Deploy on any DAPI image — no expensive transcriptomics needed at analysis time",
        "Enable retrospective analysis of archived histology slides",
    ], left_in=0.8, top_in=2.9, width_in=11.5, height_in=3.5,
       size_pt=18, line_spacing=1.5)

    add_caption(slide, "Vision: a universal DAPI classifier "
                "that works across tissues with one stain.")


def slide_pipeline_diagram(prs: Presentation) -> None:
    """Reuses the existing slide-8 illustration via inline shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_accent_bar(slide)
    add_title(slide, "DAPIDL pipeline")
    add_subtitle(slide, "From raw spatial transcriptomics → deployable DAPI model in 6 steps.",
                 top_in=1.25)

    steps = [
        ("INPUT",     "Xenium · MERSCOPE · STHELAR",   RGBColor(0x14, 0x44, 0x70)),
        ("ANNOTATE",  "11 methods → consensus",         RGBColor(0xE8, 0xA8, 0x4A)),
        ("STANDARDIZE","Cell Ontology mapping",         RGBColor(0xD0, 0xCD, 0x5C)),
        ("FILTER",    "Confidence tiers (T1/T2/T3)",    RGBColor(0x6F, 0xC2, 0x88)),
        ("EXTRACT",   "DAPI patches → LMDB",            RGBColor(0x6E, 0xB4, 0xE6)),
        ("TRAIN",     "EfficientNetV2-S → 3 heads",     RGBColor(0xA6, 0x7C, 0xD6)),
        ("DEPLOY",    "DAPI image in → cell types out", RGBColor(0x14, 0x44, 0x70)),
    ]
    box_w = 1.65
    box_h = 0.95
    spacing_x = 0.10
    total_w = len(steps) * box_w + (len(steps) - 1) * spacing_x
    start_x = (13.333 - total_w) / 2
    top = 2.6

    for i, (kicker, body, color) in enumerate(steps):
        x = start_x + i * (box_w + spacing_x)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(top),
                                       Inches(box_w), Inches(box_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0.08)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = kicker
        r.font.name = "Calibri"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run()
        r.text = body
        r.font.name = "Calibri"
        r.font.size = Pt(9)
        r.font.color.rgb = WHITE

        # Arrow chevron between boxes
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + box_w - 0.04),
                Inches(top + box_h / 2 - 0.18),
                Inches(spacing_x + 0.10),
                Inches(0.36),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    add_paragraph_block(slide, [
        "Steps 1–4 produce high-confidence training labels from gene expression.",
        "Steps 5–7 turn those labels into a DAPI-only classifier.",
        "Modular: swap any step (e.g. add new annotators, swap backbones, fuse modalities).",
    ], left_in=0.6, top_in=4.2, width_in=12.1, height_in=2.5,
       size_pt=15, line_spacing=1.3)


def slide_headline_numbers(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_accent_bar(slide)
    add_title(slide, "Three numbers to remember")
    add_subtitle(slide, "Everything else in the deck unpacks one of these.",
                 top_in=1.25)

    cards = [
        ("0.84",   "macro F1",
         "popV ensemble vs expert pathologist (Janesick et al.) on Xenium breast",
         IMM_NAVY),
        ("0.60",   "DAPI ceiling",
         "test macro F1 from DAPI-only model at coarse 4-class breast",
         EPI_RED),
        ("+0.05",  "from H&E fusion",
         "F1 gain when DAPI is fused with H&E via cross-attention",
         RGBColor(0x6F, 0x4C, 0x9C)),
    ]
    card_w = 3.95
    spacing = 0.15
    total = len(cards) * card_w + (len(cards) - 1) * spacing
    start = (13.333 - total) / 2
    top = 2.3
    h = 3.4

    for i, (number, label, blurb, accent) in enumerate(cards):
        x = start + i * (card_w + spacing)
        # Background card
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(top),
                                    Inches(card_w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = accent
        bg.line.width = Pt(2.0)

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(x), Inches(top),
                                     Inches(card_w), Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Big number
        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(top + 0.4),
                                      Inches(card_w - 0.4), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = number
        r.font.name = "Calibri"
        r.font.size = Pt(72)
        r.font.bold = True
        r.font.color.rgb = accent

        # Label
        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(top + 1.85),
                                      Inches(card_w - 0.4), Inches(0.45))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Calibri"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = NAVY_DARK

        # Description
        tb = slide.shapes.add_textbox(Inches(x + 0.3), Inches(top + 2.4),
                                      Inches(card_w - 0.6), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = blurb
        r.font.name = "Calibri"
        r.font.size = Pt(13)
        r.font.color.rgb = BODY


def slide_image(prs: Presentation, title: str, takeaway: str,
                image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_accent_bar(slide)
    add_title(slide, title)
    add_subtitle(slide, takeaway, top_in=1.25)
    if image_path.exists():
        add_image_centered(slide, image_path, top_in=1.85, max_h_in=4.95,
                           max_w_in=12.0)
    else:
        # Placeholder if the figure isn't rendered yet
        ph = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.0), Inches(2.5), Inches(9.3), Inches(3.5),
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = ACCENT
        ph.line.color.rgb = NAVY
        tf = ph.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[ {image_path.name} pending ]"
        r.font.name = "Calibri"
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = NAVY_DARK
    add_caption(slide, caption)


def slide_limitations(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_accent_bar(slide)
    add_title(slide, "Limitations & honest caveats")
    add_subtitle(slide,
                 "What this work does not (yet) deliver — listed up-front so reviewers don't have to dig.",
                 top_in=1.25)

    items = [
        ("DAPI ceiling",
         "Macro F1 plateaus around 0.55–0.65 for coarse classes. Nuclear "
         "morphology alone does not fully resolve cell identity."),
        ("Stromal & rare classes",
         "Endothelial F1 ≈ 0.45, Mast/Adipocyte F1 ≈ 0–0.20. Class imbalance "
         "and morphological similarity to neighbours are the bottleneck."),
        ("Cross-platform domain shift",
         "Models trained on one platform lose ~10–25 F1 points on another. "
         "STHELAR → Xenium transfers better than Xenium → STHELAR."),
        ("Training-label noise",
         "Spatial-transcriptomics-derived labels have ~10–15 % error rate "
         "(see Fig 4 vs Janesick GT). The model inherits this floor."),
        ("Single tissue dominance",
         "The cleanest results are on breast. Brain LOTO drops to F1 ≈ 0.10 "
         "because its class composition (93 % endothelial) differs sharply."),
    ]

    top = 2.0
    row_h = 0.95
    for i, (kicker, body) in enumerate(items):
        y = top + i * row_h
        # Kicker chip
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.6), Inches(y),
                                      Inches(2.6), Inches(0.65))
        chip.fill.solid()
        chip.fill.fore_color.rgb = NAVY_DARK
        chip.line.fill.background()
        tf = chip.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = kicker
        r.font.name = "Calibri"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE

        # Body text
        tb = slide.shapes.add_textbox(Inches(3.4), Inches(y + 0.05),
                                      Inches(9.5), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = body
        r.font.name = "Calibri"
        r.font.size = Pt(13)
        r.font.color.rgb = BODY


def slide_outlook(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_accent_bar(slide)
    add_title(slide, "Impact & outlook")
    add_subtitle(slide,
                 "If we accept the F1 ceiling, what does this enable?",
                 top_in=1.25)

    add_paragraph_block(slide, [
        "DAPI staining is already routine in most pathology / histology labs",
        "Cell types can be inferred years after tissue collection — retrospective analysis",
        "No expensive spatial transcriptomics needed at inference time",
        "Platform-agnostic: train on any spatial platform, predict on any DAPI image",
        "Modular pipeline: easy to extend with new tissues, methods, modalities",
    ], left_in=0.7, top_in=1.95, width_in=11.8, height_in=3.0,
       size_pt=17, line_spacing=1.5)

    # Highlight box
    add_box(slide, [
        ("Vision", 16, NAVY_DARK, True),
        ("A universal DAPI classifier that maps any nuclear-stained slide to cell "
         "types — at the F1 ceiling biology allows.", 15, BODY, False),
    ], left_in=0.6, top_in=5.5, width_in=12.1, height_in=1.3)


# ─── Main ──────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    slide_title(prs)

    # 2. Goal
    slide_goal(prs)

    # 3. Pipeline diagram
    slide_pipeline_diagram(prs)

    # 4. Headline numbers (new)
    slide_headline_numbers(prs)

    # 5. Annotation pipeline benchmark
    slide_image(
        prs,
        "Which annotation method should you trust?",
        "popV ensembles win on macro F1; BANKSY wins on accuracy. Choose by metric, not method.",
        FIG_DIR / "fig00_annotation_methods.png",
        "STHELAR breast s0 (574,869 cells), Cell-Ontology 7-class label. "
        "Open circles = accuracy on same data — note the F1/acc gap.",
    )

    # 6. Granularity
    slide_image(
        prs,
        "What can DAPI predict, at what granularity?",
        "F1 drops from 0.70 (4-class) to 0.48 (9-class). Per-class panel: epithelial easy, endothelial hardest.",
        FIG_DIR / "fig01_granularity.png",
        "All runs use EfficientNetV2-S, 128 px patches, identical hyperparameters.",
    )

    # 7. Patch size
    slide_image(
        prs,
        "Patch size — local context matters",
        "Larger patches consistently raise F1 across both Xenium and STHELAR.",
        FIG_DIR / "fig02_patch_size_sweep.png",
        "EffNetV2-S · Xenium 4-class · STHELAR 6-class · identical hyperparameters.",
    )

    # 8. LOTO + diversity
    slide_image(
        prs,
        "Cross-tissue generalisation — bounded by class diversity",
        "The brain paradox: when one class dominates, accuracy ≫ macro F1.",
        FIG_DIR / "fig03_loto_diversity.png",
        "Each row = one tissue held out from training; F1 computed on that tissue alone.",
    )

    # 9. Modality
    slide_image(
        prs,
        "DAPI · H&E · multimodal fusion",
        "H&E alone slightly beats DAPI; cross-attention fusion adds +5 pp F1 over DAPI.",
        FIG_DIR / "fig04_modality.png",
        "STHELAR same 9-class test split (n=189,668), EfficientNetV2-S, identical training.",
    )

    # 10. Cross-source
    slide_image(
        prs,
        "Cross-platform: STHELAR ↔ Xenium",
        "Domain shift costs 10–25 F1 points; STHELAR-trained models transfer better than Xenium-trained.",
        FIG_DIR / "fig05_cross_source.png",
        "Both panels are breast tissue. In-domain = held-out test of the SAME platform.",
    )

    # 11. Data scale
    slide_image(
        prs,
        "How much training data do you need?",
        "F1 climbs from 0.31 (5 %) to 0.42 (100 %); rare classes — endothelial, mast — benefit most.",
        FIG_DIR / "fig06_data_scale.png",
        "STHELAR breast s0 · fixed val/test (n=71,228 each) · "
        "5 fractions of the train pool · EfficientNetV2-S, identical hyperparameters.",
    )

    # 12. Limitations
    slide_limitations(prs)

    # 13. Outlook
    slide_outlook(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT} · {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
